import os
from pptx import Presentation

pptx_files = ['ĐATN-slide-original-backup.pptx', 'ĐATN-slide-clean.pptx', 'test_out.pptx', 'ĐATN-slide.pptx']

with open(r'c:\Hub\DoAn\scratch\inspect_all_pptx.txt', 'w', encoding='utf-8') as f:
    for filename in pptx_files:
        path = os.path.join(r'c:\Hub\DoAn', filename)
        if not os.path.exists(path):
            f.write(f"{filename} does not exist!\n")
            continue
        try:
            prs = Presentation(path)
            f.write(f"=== File: {filename} (Slides: {len(prs.slides)}) ===\n")
            for idx, slide in enumerate(prs.slides):
                # We are interested in Slide 7, 8, 9, 10
                if idx in [6, 7, 8, 9, 10]:
                    f.write(f"  Slide {idx+1}:\n")
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            f.write(f"    Text: {shape.text.strip()[:150]}\n")
                        if shape.shape_type == 13: # Picture
                            f.write(f"    Picture shape name: {shape.name}\n")
        except Exception as e:
            f.write(f"Error reading {filename}: {str(e)}\n")
        f.write("\n")
print("Done!")
