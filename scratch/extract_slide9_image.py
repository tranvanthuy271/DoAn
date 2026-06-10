from pptx import Presentation

prs = Presentation(r'c:\Hub\DoAn\ĐATN-slide-clean.pptx')
slide = prs.slides[8] # Slide 9 (index 8)
print("Slide title:", slide.shapes.title.text if slide.shapes.title else "No title")

for idx, shape in enumerate(slide.shapes):
    if shape.shape_type == 13: # Picture
        image = shape.image
        ext = image.ext
        print(f"  Image ext: {ext}, size: {len(image.blob)} bytes")
        with open(f"c:\\Hub\\DoAn\\scratch\\slide9_extracted_img.{ext}", "wb") as img_file:
            img_file.write(image.blob)
        print("  Saved extracted image to c:\\Hub\\DoAn\\scratch\\slide9_extracted_img." + ext)
