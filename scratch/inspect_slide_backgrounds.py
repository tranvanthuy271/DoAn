import pptx

prs = pptx.Presentation(r'c:\Hub\DoAn\ĐATN-slide.pptx')
slide = prs.slides[8] # Slide 9

print("Background fill type:", slide.background.fill.type if slide.background else "No background object")
print("Shapes count:", len(slide.shapes))
for idx, shape in enumerate(slide.shapes):
    print(f"Shape {idx}: name='{shape.name}', type={shape.shape_type}")

# Let's inspect slide xml to see where the image or text is
import xml.etree.ElementTree as ET
slide_xml = slide._element.xml
print("\nSlide XML (first 1000 chars):")
print(slide_xml[:1000])
