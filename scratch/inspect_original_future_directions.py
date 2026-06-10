import pptx

prs = pptx.Presentation(r'c:\Hub\DoAn\ĐATN-slide-clean.pptx')
slide = prs.slides[28] # Slide 29 (0-indexed 28)
with open(r'c:\Hub\DoAn\scratch\original_future_directions.txt', 'w', encoding='utf-8') as f:
    f.write("=== Slide 29 Title ===\n")
    f.write((slide.shapes.title.text if slide.shapes.title else "No Title") + "\n\n")
    for s_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame and shape.text.strip():
            f.write(f"Shape {s_idx}:\n")
            f.write(shape.text.strip() + "\n\n")
print("Done!")
