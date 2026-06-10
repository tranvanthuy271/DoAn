import pptx

prs = pptx.Presentation(r'c:\Hub\DoAn\ĐATN-slide.pptx')
slide = prs.slides[8] # Slide 9 (0-indexed 8)

print("=== Slide 9 Details ===")
for idx, shape in enumerate(slide.shapes):
    print(f"Shape {idx}: name='{shape.name}', type={shape.shape_type}")
    if shape.has_text_frame and shape.text.strip():
        print(f"  Text: {shape.text.strip()[:200]}")
    if shape.shape_type == 13: # Picture
        print(f"  Image: size={len(shape.image.blob)} bytes")
