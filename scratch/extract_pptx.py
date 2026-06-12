import sys
import os
import subprocess

# Reconfigure stdout to use utf-8
try:
    sys.stdout.reconfigure(encoding='utf-8')
except AttributeError:
    pass

try:
    import pptx
except ImportError:
    print("python-pptx not found, installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation

def extract_text_from_pptx(file_path, output_path):
    if not os.path.exists(file_path):
        print(f"Error: File {file_path} not found.")
        return

    prs = Presentation(file_path)
    print(f"Total Slides: {len(prs.slides)}")
    
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"Total Slides: {len(prs.slides)}\n\n")
        
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            
            # Extract title if exists
            title = ""
            if slide.shapes.title:
                try:
                    title = slide.shapes.title.text.strip()
                    f.write(f"Title: {title}\n")
                except Exception as e:
                    f.write(f"Title Error: {str(e)}\n")
            
            # Extract text from shapes
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Avoid printing the title again if we already got it
                    if slide.shapes.title and shape == slide.shapes.title:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = paragraph.text.strip()
                        if p_text:
                            texts.append(p_text)
            
            if texts:
                f.write("Content:\n")
                for txt in texts:
                    f.write(f"  - {txt}\n")
            else:
                f.write("Content: (No text found or image/diagram only)\n")
            f.write("\n")
            
    print(f"Successfully extracted slides to {output_path}")

if __name__ == "__main__":
    pptx_path = "ĐATN-TVT.pptx"
    output_path = "scratch/extracted_slides.txt"
    extract_text_from_pptx(pptx_path, output_path)
