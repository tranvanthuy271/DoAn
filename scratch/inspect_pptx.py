import collections
import collections.abc
from pptx import Presentation

prs = Presentation("ĐATN-TVT.pptx")
print(f"Total slides: {len(prs.slides)}")

with open("scratch/pptx_content.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n=================== SLIDE {i+1} ===================\n")
        # Extract title if exists
        if slide.shapes.title:
            f.write(f"TITLE: {slide.shapes.title.text}\n")
        
        # Extract other text boxes
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                f.write(f"TEXT:\n")
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        f.write(f"  {paragraph.text}\n")
                        
print("Slide extraction completed successfully.")
