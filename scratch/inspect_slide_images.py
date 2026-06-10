import sys
from pptx import Presentation

prs = Presentation(r'c:\Hub\DoAn\ĐATN-slide-original-backup.pptx')
print("Total slides:", len(prs.slides))

with open(r'c:\Hub\DoAn\scratch\pptx_images_info.txt', 'w', encoding='utf-8') as f:
    for idx, slide in enumerate(prs.slides):
        f.write(f"Slide {idx+1}:\n")
        # Find images on this slide
        img_idx = 1
        for shape in slide.shapes:
            if shape.name.startswith('Picture') or shape.shape_type == 13: # 13 is MSO_SHAPE_TYPE.PICTURE
                f.write(f"  Shape {shape.name}: image shape_type={shape.shape_type}\n")
                img_idx += 1
            elif shape.has_text_frame:
                f.write(f"  Text: {shape.text[:100].strip()}\n")
        f.write("\n")
print("Done inspecting shapes!")
