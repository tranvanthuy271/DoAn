import pptx

prs = pptx.Presentation(r'c:\Hub\DoAn\ĐATN-slide-clean.pptx')
with open(r'c:\Hub\DoAn\scratch\original_final_slides.txt', 'w', encoding='utf-8') as f:
    for idx in range(29, len(prs.slides)):
        slide = prs.slides[idx]
        f.write(f"=== Slide {idx+1} ===\n")
        for s_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame and shape.text.strip():
                f.write(f"  Shape {s_idx} ({shape.name}) text:\n")
                f.write("    " + "\n    ".join(shape.text.split('\n')) + "\n")
            if shape.shape_type == 13:
                f.write(f"  Shape {s_idx} ({shape.name}) is image\n")
        f.write("\n")
print("Done!")
