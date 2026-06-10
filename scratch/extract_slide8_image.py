from pptx import Presentation

prs = Presentation(r'c:\Hub\DoAn\ĐATN-slide-clean.pptx')
slide = prs.slides[7] # Slide 8 (index 7)
print("Slide title:", slide.shapes.title.text if slide.shapes.title else "No title")

for idx, shape in enumerate(slide.shapes):
    print(f"Shape {idx}: name={shape.name}, type={shape.shape_type}")
    if shape.shape_type == 13: # Picture
        image = shape.image
        # Save image bytes to check format and size
        ext = image.ext
        print(f"  Image ext: {ext}, size: {len(image.blob)} bytes")
        with open(f"c:\\Hub\\DoAn\\scratch\\slide8_extracted_img.{ext}", "wb") as img_file:
            img_file.write(image.blob)
        print("  Saved extracted image to c:\\Hub\\DoAn\\scratch\\slide8_extracted_img." + ext)
