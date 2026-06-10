import pptx

prs = pptx.Presentation(r'c:\Hub\DoAn\ĐATN-slide.pptx')
with open(r'c:\Hub\DoAn\scratch\pptx_structure.txt', 'w', encoding='utf-8') as f:
    f.write(f"Number of slides: {len(prs.slides)}\n\n")
    for idx, slide in enumerate(prs.slides):
        f.write(f"=== Slide {idx+1} ===\n")
        # Check slide title if exists
        if slide.shapes.title:
            f.write(f"Title: {slide.shapes.title.text}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        f.write(f"  {paragraph.text}\n")
        f.write("\n")
print("Successfully dumped PPTX structure!")
