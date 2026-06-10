import os
import hashlib
from pptx import Presentation

prs = Presentation(r'c:\Hub\DoAn\ĐATN-slide-clean.pptx')

# MD5 hashes of all extracted images
extracted_hashes = {}
for f in os.listdir(r'c:\Hub\DoAn\extracted_images'):
    p = os.path.join(r'c:\Hub\DoAn\extracted_images', f)
    if os.path.isfile(p):
        with open(p, 'rb') as file_obj:
            h = hashlib.md5(file_obj.read()).hexdigest()
            extracted_hashes[h] = f

with open(r'c:\Hub\DoAn\scratch\all_clean_slide_images.txt', 'w', encoding='utf-8') as f_out:
    for idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text.strip() else "[No Title]"
        f_out.write(f"Slide {idx+1}: {title}\n")
        for s_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13: # Picture
                image = shape.image
                h = hashlib.md5(image.blob).hexdigest()
                mapped_name = extracted_hashes.get(h, "Unknown image")
                f_out.write(f"  Shape {s_idx} ({shape.name}): {mapped_name} (ext: {image.ext}, size: {len(image.blob)})\n")
        f_out.write("\n")
print("Done mapping all clean slide images!")
