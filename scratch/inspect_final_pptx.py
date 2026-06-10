import pptx

prs = pptx.Presentation(r'c:\Hub\DoAn\ĐATN-slide.pptx')
print("Total slides:", len(prs.slides))

with open(r'c:\Hub\DoAn\scratch\inspect_final_pptx.txt', 'w', encoding='utf-8') as f:
    for idx, slide in enumerate(prs.slides):
        f.write(f"Slide {idx+1}:\n")
        for s_idx, shape in enumerate(slide.shapes):
            f.write(f"  Shape {s_idx}: name={shape.name}, type={shape.shape_type}\n")
            if shape.has_text_frame and shape.text.strip():
                f.write(f"    Text: {shape.text.strip()[:100]}\n")
        f.write("\n")
print("Done!")
