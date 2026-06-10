import os
import sys
import docx
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import pptx
from pptx.util import Inches as PptxInches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

# Configure stdout to handle UTF-8 output
sys.stdout.reconfigure(encoding='utf-8')

# Output paths
docx_out_path = r"c:\Hub\DoAn\KichBanSlide.docx"
pptx_out_path = r"c:\Hub\DoAn\ĐATN-slide-clean.pptx"
md_out_path = r"c:\Hub\DoAn\docs\ĐATN-slide.md"
img_dir = r"c:\Hub\DoAn\extracted_images"

# Slide data definition
slides_data = [
    {
        "slide_num": 1,
        "title": "PHÁT TRIỂN TRÒ CHƠI MUTANTS ARENA\nVỚI HỆ THỐNG TIẾN HÓA GENE BẰNG UNITY",
        "subtitle": "ĐỒ ÁN TỐT NGHIỆP ĐẠI HỌC",
        "type": "cover",
        "content": [
            "Học viện Kỹ thuật Mật mã – Khoa Công nghệ thông tin",
            "Sinh viên thực hiện: Trần Văn Thủy – Lớp: CT6 – MSV: CT060439",
            "Người hướng dẫn: TS. Nguyễn Đức Hiếu",
            "Năm thực hiện: 2026"
        ],
        "images": []
    },
    {
        "slide_num": 2,
        "title": "Đặt vấn đề & Mục tiêu Đề tài",
        "subtitle": "Đặt vấn đề",
        "type": "standard",
        "content": [
            "• Bối cảnh: Trò chơi nhập vai hành động 2D (Action RPG) đang bùng nổ mạnh mẽ, đặc biệt ở phân khúc game độc lập (indie).",
            "• Thực trạng: Thị trường game Việt Nam chủ yếu là game nhập khẩu; chưa có nhiều sản phẩm tự sản xuất có chiều sâu về gameplay và tích hợp mạng multiplayer phức tạp.",
            "• Thách thức kỹ thuật: Xây dựng hệ thống di chuyển platformer mượt mà (dash, coyote time, jump buffer), phát triển cơ chế tiến hóa nhân vật độc đáo và chống hack/cheat vị trí/tài nguyên trên môi trường mạng trực tuyến.",
            "• Mục tiêu đề tài: Phát triển trò chơi multiplayer Mutants Arena góc nhìn 2D Platformer, tích hợp cơ chế tiến hóa Gene Ngũ Hành khắc chế linh hoạt, triển khai Dedicated Server theo kiến trúc Server-Authoritative để chống gian lận."
        ],
        "images": ["image_unnamed_7.png"] # Figure 2.1 Sơ đồ kiến trúc tổng thể
    },
    {
        "slide_num": 3,
        "title": "Nội dung báo cáo",
        "subtitle": "Mục lục",
        "type": "standard",
        "content": [
            "• Chương 1: Tổng quan về đề tài và cơ sở lý thuyết",
            "  - Giới thiệu trò chơi, các cơ chế di chuyển, FSM quái vật, kiến trúc Server-Authoritative, dự đoán client, đồng bộ qua Netcode for GameObjects và SignalR.",
            "• Chương 2: Phân tích và thiết kế hệ thống",
            "  - Sơ đồ kiến trúc tổng thể, đặc tả Use Case tổng quát và chi tiết (Chiến đấu, Tiến hóa Gene, Phó bản), biểu đồ tuần tự (Sequence Diagram), sơ đồ cơ sở dữ liệu vật lý (ERD).",
            "• Chương 3: Triển khai chi tiết và kết quả thực nghiệm",
            "  - Môi trường triển khai VPS Linux bằng Docker Compose, cơ chế bảo mật nhiều lớp (Zone API Key, Connection Approval) và hình ảnh kết quả thực tế các tính năng trò chơi."
        ],
        "images": []
    },
    {
        "slide_num": 4,
        "title": "Chương 1: Tổng quan và Cơ sở lý thuyết",
        "subtitle": "Chương 1 giới thiệu",
        "type": "divider",
        "content": [
            "Chương 1 trình bày tổng quan về dự án trò chơi Mutants Arena và cơ sở lý thuyết nền tảng để xây dựng game di chuyển platformer thời gian thực đa người chơi."
        ],
        "images": []
    },
    {
        "slide_num": 5,
        "title": "Cơ chế di chuyển & Hoạt ảnh nhân vật",
        "subtitle": "Animator Controller & Movement feel",
        "type": "standard",
        "content": [
            "• Animator Controller: Đồ thị trạng thái quản lý chuyển đổi mượt mà giữa các animation (Idle, Run, Jump, Fall, Dash, Attack, Die) dựa trên Exit Time = 0 và Transition Duration ngắn (0.05 - 0.1 giây).",
            "• Coyote Time: Cho phép người chơi nhảy khi đã rơi khỏi mép vực trong thời gian ngắn (0.1 giây) để giảm ức chế điều khiển.",
            "• Jump Buffer: Lưu lệnh nhảy của người chơi khi chưa chạm đất (khoảng 0.15 giây trước khi chạm) và tự kích hoạt ngay khi tiếp đất.",
            "• Dash Invincibility (I-Frames): Trạng thái bất tử ngắn hạn khi Dash, giúp nhân vật tránh sát thương hoàn toàn, tăng tính thử thách và kỹ năng trong chiến đấu."
        ],
        "images": ["image_unnamed_1.png"] # Figure 1.1 Animator Controller
    },
    {
        "slide_num": 6,
        "title": "Cơ chế Tiến hóa Gene & Tương khắc Ngũ Hành",
        "subtitle": "Gene Evolution & Element",
        "type": "standard",
        "content": [
            "• Tiến hóa Gene: Hệ thống thay thế Class truyền thống. Người chơi nâng cấp Gene từ Tier 1 (Sơ cấp) lên Tier 5 (Tối thượng) để nhận kỹ năng chủ động, bị động và chỉ số thuộc tính.",
            "• Cấu trúc Multi-Gene: Cho phép khảm tối đa 1 Gene chính (100% hiệu quả chỉ số) và 2 Gene phụ (30% hiệu quả chỉ số) tạo sự linh hoạt trong build đồ.",
            "• Hybrid Fusion: Dung hợp 2 Gene hệ khác nhau ở Tier 5 để tạo class lai (Hybrid) đặc biệt.",
            "• Vòng tương khắc Ngũ Hành: Kim ➔ Mộc ➔ Thủy ➔ Hỏa ➔ Thổ ➔ Kim. Tấn công khắc hệ nhân x1.5 sát thương; bị khắc hệ nhân x0.75 sát thương; hệ trung lập nhân x1.0.",
            "• Hệ thứ 6 Phong (Wind): Hệ trung lập đặc biệt, không khắc và không bị khắc bởi 5 hệ Ngũ Hành, giúp cân bằng lối chơi."
        ],
        "images": ["image_unnamed_52.png"] # Figure 3.7 Chọn hệ nguyên tố
    },
    {
        "slide_num": 7,
        "title": "Hệ thống Quái vật AI (FSM)",
        "subtitle": "AI FSM quái vật",
        "type": "standard",
        "content": [
            "• Quái vật AI: Sử dụng Finite State Machine (FSM) để mô hình hóa hành vi trên Dedicated Server nhằm tối ưu hóa CPU.",
            "• Các trạng thái chính của quái vật:",
            "  - Idle / Patrol: Đi tuần tra quanh vùng xuất hiện (Spawn Anchor Point).",
            "  - Chase: Đuổi theo người chơi gần nhất lọt vào tầm phát hiện (Detection Range).",
            "  - Attack: Thực hiện đòn đánh khi người chơi trong tầm tấn công (Attack Range) và kỹ năng hết thời gian hồi (cooldown).",
            "  - Hit / Dead: Khựng lại nhận hiệu ứng hoặc tự giải phóng tài nguyên và kích hoạt drop rate khi chết."
        ],
        "images": ["image_unnamed_2.png"] # Figure 1.2 Sơ đồ FSM quái vật
    },
    {
        "slide_num": 8,
        "title": "Kiến trúc mạng Server-Authoritative",
        "subtitle": "Server-Authoritative Architecture",
        "type": "standard",
        "content": [
            "• Nguyên lý cốt lõi: Dedicated Server chạy độc lập trên VPS là nguồn thông tin duy nhất và đáng tin cậy (Single Source of Truth). Client chỉ gửi input điều khiển và nhận gói tin đồng bộ trạng thái.",
            "• Luồng xử lý kỹ năng và di chuyển:",
            "  1. Client gửi yêu cầu ServerRpc (SkillInput/MoveInput) lên Server.",
            "  2. Server thực thi kiểm tra tính hợp lệ trên Server (MP, cooldown, vị trí va chạm).",
            "  3. Server xác nhận kết quả, cập nhật HP/vị trí và broadcast gói tin ClientRpc về cho tất cả clients trong vùng hiển thị.",
            "• Chống gian lận: Client không thể tự ý thay đổi HP, vàng, sát thương hoặc bay nhảy tự do vì Server liên tục xác thực tất cả logic vật lý và dữ liệu."
        ],
        "images": ["image_unnamed_3.png"] # Figure 1.3 Luồng Server-Authoritative
    },
    {
        "slide_num": 9,
        "title": "Client-side Prediction & Server Reconciliation",
        "subtitle": "Đồng bộ hóa di chuyển chống trễ mạng",
        "type": "standard",
        "content": [
            "• Bài toán trễ mạng (Latency): Nếu di chuyển phải đợi Server phản hồi (RTT ~100ms), người chơi sẽ cảm nhận độ khựng trễ cực kỳ khó chịu.",
            "• Client-side Prediction: Client tự áp dụng input di chuyển của người chơi lập tức lên nhân vật cục bộ để đảm bảo cảm giác điều khiển mượt mà tức thì.",
            "• Server Reconciliation: Server gửi lại vị trí chính xác (State) về. Client so sánh vị trí cục bộ ở thời điểm tương ứng. Nếu sai lệch vượt quá ngưỡng cho phép (threshold), Client tự động điều chỉnh mượt mà (reconcile) về vị trí đúng của Server để loại bỏ hiện tượng giật giật (rubber banding)."
        ],
        "images": ["image_unnamed_47.png"] # Figure 3.3 Luồng kiểm duyệt NGO Dedicated Server
    },
    {
        "slide_num": 10,
        "title": "Đồng bộ qua Unity NGO & SignalR",
        "subtitle": "Mô hình truyền tải dữ liệu game",
        "type": "standard",
        "content": [
            "• Unity NGO (Netcode for GameObjects):",
            "  - Chạy trên giao thức UDP (unity-transport) tối ưu cho truyền gói tin tần suất cao.",
            "  - Đồng bộ di chuyển (NetworkTransform), hoạt ảnh nhân vật, vị trí quái vật và các đòn đánh qua NetworkVariable và RPC.",
            "• SignalR (WebSockets/TCP):",
            "  - Chạy trên giao thức WebSockets đảm bảo tin cậy cao.",
            "  - Quản lý các tính năng meta-game: Lobby sảnh chờ, kết nối tổ đội (Party), gửi tin nhắn kênh chat thế giới/tổ đội, quản lý bạn bè thời gian thực.",
            "• Phân chia trách nhiệm: NGO cho gameplay realtime trễ thấp; SignalR cho dữ liệu xã hội và meta-game cần độ tin cậy cao."
        ],
        "images": ["image_unnamed_45.png"] # Figure 3.2 Luồng xác thực nội bộ
    },
    {
        "slide_num": 11,
        "title": "Kiến trúc Zone-based Server",
        "subtitle": "Phân khu máy chủ giảm tải băng thông",
        "type": "standard",
        "content": [
            "• Phân chia Zone: Bản đồ game chia thành các phân khu bản đồ độc lập (Zone). Một Client chỉ nhận gói tin đồng bộ từ những người chơi và quái vật trong cùng một Zone, giúp giảm tải CPU và băng thông Server cực lớn.",
            "• ZoneRoomRegistry: Tiến trình Server quản lý tập trung danh sách địa chỉ IP/Port của các Zone và các instance phó bản đang hoạt động.",
            "• Instance (Phó bản động): Khi tổ đội xin đi phó bản, Server tạo động một Instance bản đồ riêng tư, phân phối người chơi vào đó và tự động giải phóng tài nguyên sau khi phó bản kết thúc."
        ],
        "images": ["image_unnamed_6.png"] # Figure 1.6 Kiến trúc Docker Compose
    },
    {
        "slide_num": 12,
        "title": "Chương 2: Phân tích và Thiết kế hệ thống",
        "subtitle": "Chương 2 giới thiệu",
        "type": "divider",
        "content": [
            "Chương 2 tập trung phân tích các yêu cầu chức năng hệ thống Mutants Arena và thực hiện thiết kế chi tiết thông qua các biểu đồ ca sử dụng (Use Case), biểu đồ tuần tự (Sequence Diagram) và sơ đồ cơ sở dữ liệu (ERD)."
        ],
        "images": []
    },
    {
        "slide_num": 13,
        "title": "Mô hình kiến trúc tổng thể hệ thống",
        "subtitle": "Kiến trúc 3 lớp",
        "type": "standard",
        "content": [
            "• Tầng Client: Unity Client xây dựng trên ngôn ngữ C#, xử lý render hình ảnh, hoạt ảnh nhân vật, chạy Client-side prediction và gửi đầu vào điều khiển mạng.",
            "• Tầng Server:",
            "  - Dedicated Server (Unity Headless Build) chạy logic gameplay thời gian thực.",
            "  - API & SignalR Server (ASP.NET Core 8.0) quản lý xác thực tài khoản, lưu dữ liệu nhân vật, quản lý kênh chat và tổ đội.",
            "• Tầng Database: MySQL Database lưu trữ an toàn dữ liệu tài khoản, chỉ số nhân vật, cấu hình Gene, túi đồ qua EF Core."
        ],
        "images": ["image_unnamed_7.png"] # Figure 2.1 Sơ đồ kiến trúc tổng thể
    },
    {
        "slide_num": 14,
        "title": "Biểu đồ Use Case mức tổng quát",
        "subtitle": "Tác nhân và Ca sử dụng",
        "type": "standard",
        "content": [
            "• Tác nhân Player: Người chơi trực tiếp tham gia vào game thông qua client để thực hiện đăng nhập, di chuyển, chiến đấu, tiến hóa Gene, dung hợp, chat, tạo tổ đội, đi phó bản và nâng cấp trang bị.",
            "• Tác nhân Dedicated Server: Game Server đóng vai trò tự động xử lý, đồng bộ và xác thực các hành động di chuyển, chiến đấu, spawn quái vật, hoàn tất và phát thưởng phó bản.",
            "• Tác nhân API Server: Xử lý các nghiệp vụ đăng ký/đăng nhập, lưu dữ liệu nhân vật, cập nhật giao dịch túi đồ và bảng xếp hạng."
        ],
        "images": ["image_unnamed_8.png"] # Figure 2.2 Biểu đồ ca sử dụng mức tổng quát
    },
    {
        "slide_num": 15,
        "title": "Use Case chi tiết: Chiến đấu & Sử dụng kỹ năng",
        "subtitle": "Use Case Combat",
        "type": "standard",
        "content": [
            "• Đăng ký đầu vào: Người chơi nhấn phím kỹ năng (Q, W, E, R). Client kiểm tra cooldown cục bộ và gửi SkillInputServerRpc lên Server.",
            "• Xác thực trên Server: Dedicated Server nhận yêu cầu, kiểm tra lượng MP hiện tại của nhân vật, trạng thái khống chế (stun, freeze) và cooldown thực tế.",
            "• Tính toán sát thương: Server thực hiện va chạm vật lý (OverlapCircleAll) xác định quái vật trúng chiêu. Tính toán lượng sát thương dựa trên chỉ số nhân vật và ma trận khắc chế nguyên tố Ngũ Hành.",
            "• Phản hồi: Trừ MP nhân vật, trừ HP quái vật, gửi gói tin ClientRpc đồng bộ hiệu ứng đòn đánh và lượng HP thay đổi về tất cả clients trong zone."
        ],
        "images": ["image_unnamed_12.png"] # Figure 2.6 Biểu đồ ca sử dụng cho mô-đun Chiến đấu
    },
    {
        "slide_num": 16,
        "title": "Use Case chi tiết: Tiến hóa Gene Ngũ Hành",
        "subtitle": "Use Case Gene Evolution",
        "type": "standard",
        "content": [
            "• Yêu cầu nâng cấp: Người chơi mở giao diện Gene, chọn nút Gene muốn nâng cấp (Tier 1 - 5).",
            "• Kiểm tra điều kiện: Hệ thống kiểm tra cấp độ nhân vật hiện tại và số lượng tài nguyên yêu cầu (Gene Fragment, Gene Core, Gene Essence).",
            "• Xử lý nâng cấp: Client gửi request lên API Server. API Server xác thực, trừ tài nguyên trong DB, lưu trạng thái Gene mới và trả về thành công.",
            "• Đồng bộ chỉ số: API Server đồng bộ dữ liệu mới về Client. Client kích hoạt hiệu ứng nâng cấp và cập nhật vĩnh viễn thuộc tính nhân vật (ATK, HP, DEF, CRIT) hoặc mở khóa kỹ năng tương ứng."
        ],
        "images": ["image_unnamed_15.png"] # Figure 2.9 Biểu đồ ca sử dụng cho mô-đun Phát triển Gene
    },
    {
        "slide_num": 17,
        "title": "Use Case chi tiết: Tham gia phó bản Wave",
        "subtitle": "Use Case Dungeon Wave",
        "type": "standard",
        "content": [
            "• Tạo tổ đội: Người chơi mời bạn bè hoặc đi đơn lẻ đến gặp NPC phó bản. Gửi yêu cầu đi phó bản.",
            "• Khởi tạo Instance: Server kiểm tra cấp độ tổ đội, sau đó tạo một Instance bản đồ phó bản riêng tư trên Dedicated Server.",
            "• Trận chiến theo đợt (Waves): Phó bản gồm nhiều đợt quái spawn tăng dần độ khó. Người chơi phải tiêu diệt sạch đợt hiện tại để đợt tiếp theo xuất hiện.",
            "• Hoàn thành phó bản: Tiêu diệt Boss ở đợt cuối cùng. Server tính toán thành tích, lưu phần thưởng (Exp, Gold, Gene Core) vào DB và chuyển người chơi về lại Lobby."
        ],
        "images": ["image_unnamed_21.png"] # Figure 2.15 Biểu đồ ca sử dụng cho mô-đun Tham gia phó bản
    },
    {
        "slide_num": 18,
        "title": "Biểu đồ tuần tự: Chiến đấu & Sử dụng kỹ năng",
        "subtitle": "Sequence Diagram Combat",
        "type": "standard",
        "content": [
            "• Mô tả quy trình giao tiếp mạng thời gian thực khi người chơi tung đòn:",
            "  1. Player Client gửi Skill Input qua mạng UDP (NGO ServerRpc).",
            "  2. Dedicated Server chạy logic kiểm tra khoảng cách, HP và MP nhân vật.",
            "  3. Dedicated Server tính toán lượng sát thương và trừ máu quái vật vật lý.",
            "  4. Dedicated Server gọi API Server qua mạng nội bộ để lưu trữ các thay đổi trạng thái nếu có.",
            "  5. Dedicated Server broadcast ClientRpc gửi VFX/SFX và lượng HP quái vật giảm đi về cho tất cả clients trong vùng hiển thị."
        ],
        "images": ["image_unnamed_25.png"] # Figure 2.19 Biểu đồ tuần tự đặc tả ca sử dụng Chiến đấu
    },
    {
        "slide_num": 19,
        "title": "Biểu đồ tuần tự: Dung hợp Hybrid Gene",
        "subtitle": "Sequence Diagram Hybrid Fusion",
        "type": "standard",
        "content": [
            "• Mô tả quy trình giao tiếp khi dung hợp Gene Tier 5 tạo class lai:",
            "  1. Người chơi tương tác với UI và chọn 2 Gene hệ khác nhau đạt cấp tối đa (ví dụ: Hỏa và Lôi). Gửi yêu cầu Dung hợp lên API Server.",
            "  2. API Server truy vấn database kiểm tra tài nguyên tiêu hao và trạng thái Gene.",
            "  3. API Server cập nhật DB: Hủy bỏ 2 Gene cũ, thêm mới Hybrid Gene cấp 1, trừ tài nguyên.",
            "  4. API Server trả về kết quả thành công cho Client. Client hiển thị hoạt ảnh dung hợp thành công và mở khóa bộ kỹ năng Hybrid đặc biệt."
        ],
        "images": ["image_unnamed_28.png"] # Figure 2.22 Biểu đồ tuần tự đặc tả ca sử dụng Dung hợp Gene
    },
    {
        "slide_num": 20,
        "title": "Sơ đồ Cơ sở dữ liệu vật lý (ERD)",
        "subtitle": "Database Schema & JSON Column",
        "type": "standard",
        "content": [
            "• Sơ đồ ERD: Gồm các bảng thực thể liên kết chặt chẽ: Users (tài khoản), Player_Data (chỉ số nhân vật), Inventory (vật phẩm), Gene_Data (dữ liệu Gene).",
            "• Thiết kế JSON Column linh hoạt trong Game Database:",
            "  - Các dữ liệu có cấu trúc động như: danh sách chỉ số phụ trang bị cường hóa ngẫu nhiên, trạng thái các nhiệm vụ đang nhận, danh sách kỹ năng gán phím tắt nhanh được lưu trữ trực tiếp dưới dạng JSON trong một cột duy nhất.",
            "  - Lợi ích: Tránh việc tạo quá nhiều bảng liên kết 1-nhiều phức tạp, giảm thiểu số lượng truy vấn JOIN lớn, giúp tăng tốc độ đọc/ghi dữ liệu nhân vật lên đến 40%."
        ],
        "images": ["Hình 2.24. Sơ đồ cơ sở dữ liệu của hệ thống Mutants Arena.png"]
    },
    {
        "slide_num": 21,
        "title": "Chương 3: Triển khai chi tiết và Kết quả thực nghiệm",
        "subtitle": "Chương 3 giới thiệu",
        "type": "divider",
        "content": [
            "Chương 3 trình bày chi tiết về môi trường triển khai thực tế trên VPS Linux sử dụng Docker Compose, cơ chế bảo mật mạng nhiều lớp và hình ảnh các kết quả giao diện các tính năng đã phát triển."
        ],
        "images": []
    },
    {
        "slide_num": 22,
        "title": "Kiến trúc triển khai Docker Compose",
        "subtitle": "Docker compose VPS deployment",
        "type": "standard",
        "content": [
            "• Kiến trúc đóng gói Container: Đóng gói toàn bộ các dịch vụ hệ thống thành 3 Container độc lập:",
            "  - Container 1: MySQL Database Server (expose cổng database nội bộ cô lập).",
            "  - Container 2: Web API & SignalR Server (ASP.NET Core 8.0, expose cổng HTTP 5000).",
            "  - Container 3: Dedicated Server (Unity headless build, expose cổng UDP 7777).",
            "• Docker Network: Các Container giao tiếp tốc độ cao qua Docker network nội bộ, database hoàn toàn cô lập khỏi internet nhằm bảo mật dữ liệu nhân vật."
        ],
        "images": ["image_unnamed_49.png"] # Figure 3.4 Sơ đồ kiến trúc triển khai thực tế
    },
    {
        "slide_num": 23,
        "title": "Cơ chế bảo mật hệ thống nhiều lớp",
        "subtitle": "Security mechanisms",
        "type": "standard",
        "content": [
            "• Lớp REST API: Mã hóa mật khẩu người chơi (BCrypt), xác thực an toàn qua token JWT kèm thời gian hết hạn ngắn.",
            "• Lớp Zone Server Communication (Zone API Key):",
            "  - Khi Dedicated Server gọi REST API để cập nhật dữ liệu (vàng, exp, đồ drop) của người chơi, request phải đính kèm Zone API Key bí mật cấu hình trong Docker environment.",
            "  - Ngăn chặn kẻ xấu giả lập request API ngoài luồng để tự buff vật phẩm.",
            "• Lớp Dedicated Server (Connection Approval): Khi kết nối mạng UDP, client phải gửi token JWT để Dedicated Server xác thực danh tính với API Server trước khi chấp nhận cho vào game."
        ],
        "images": ["image_unnamed_41.png"] # Figure 3.1 Sơ đồ bảo mật nhiều lớp
    },
    {
        "slide_num": 24,
        "title": "Kết quả: Giao diện Đăng ký & Đăng nhập",
        "subtitle": "Kết quả giao diện Đăng nhập",
        "type": "standard",
        "content": [
            "• Màn hình đăng nhập thiết kế theo chủ đề Sci-Fi huyền bí phù hợp cốt truyện Mutants Arena, hiển thị logo động.",
            "• Tích hợp logic xử lý lỗi: Hiển thị thông báo trực quan trên UI khi tài khoản không tồn tại, sai mật khẩu, hoặc mất kết nối máy chủ.",
            "• Xác thực phía Client: Kiểm tra định dạng đầu vào trước khi gửi request để giảm tải truy vấn vô ích lên API Server."
        ],
        "images": ["image_unnamed_50.png"] # Figure 3.5 Giao diện đăng nhập
    },
    {
        "slide_num": 25,
        "title": "Kết quả: Giao diện Chọn nguyên tố & Sảnh chính",
        "subtitle": "Kết quả giao diện Sảnh chính",
        "type": "standard",
        "content": [
            "• Sảnh chọn hệ nguyên tố: Người chơi lựa chọn 1 trong 5 nguyên tố cơ bản khi khởi tạo nhân vật, kích hoạt bộ kỹ năng sơ cấp đặc trưng.",
            "• Giao diện sảnh chính (Lobby): Cho phép hiển thị mô hình nhân vật 2D, thanh trạng thái (HP, MP, Level, Exp), danh sách nhiệm vụ đang nhận, bảng chat đa kênh và menu truy cập nhanh các tính năng xã hội."
        ],
        "images": ["image_unnamed_52.png"] # Figure 3.7 Chọn hệ nguyên tố
    },
    {
        "slide_num": 26,
        "title": "Kết quả: Giao diện Nâng cấp Gene chính/phụ",
        "subtitle": "Kết quả giao diện Gene",
        "type": "standard",
        "content": [
            "• Giao diện Tiến hóa Gene chính: Hiển thị cây kỹ năng phân cấp trực quan theo các Tier 1 - Tier 5. Các nút Gene đã mở khóa sẽ sáng lên kèm hiệu ứng nguyên tố tương ứng.",
            "• Giao diện Nâng cấp Gene phụ: Cho phép khảm tối đa 2 Gene phụ bổ trợ chỉ số. Hiển thị mô tả rõ ràng lượng chỉ số gia tăng (ví dụ: +15 ATK, +5% Crit Rate)."
        ],
        "images": ["image_unnamed_59.png", "image_unnamed_61.png"] # Figure 3.14 & 3.16 Nâng cấp Gene chính/phụ
    },
    {
        "slide_num": 27,
        "title": "Kết quả: Giao diện Dung hợp Hybrid Gene",
        "subtitle": "Kết quả giao diện dung hợp",
        "type": "standard",
        "content": [
            "• Giao diện dung hợp: Người chơi đưa vào 2 Gene nguyên tố Tier 5 khác nhau, hệ thống tự động tính toán và hiển thị tỷ lệ dung hợp thành công.",
            "• Bảng tóm tắt nhân vật sau dung hợp: Hiển thị chi tiết lớp nhân vật Hybrid mới được mở khóa, màu sắc hào quang (Aura) đặc trưng quanh nhân vật và các chỉ số sức mạnh vượt trội."
        ],
        "images": ["image_unnamed_62.png"] # Figure 3.17 Giao diện dung hợp Hybrid
    },
    {
        "slide_num": 28,
        "title": "Kết quả: Giao diện Chỉ số, Trang bị & Kỹ năng",
        "subtitle": "Kết quả giao diện trang bị kỹ năng",
        "type": "standard",
        "content": [
            "• Tab Chỉ số & Trang bị: Hiển thị mô hình 2D nhân vật trực quan, các ô khảm trang bị (Vũ khí, Giáp, Dây chuyền, Nhẫn) và các chỉ số cơ bản của nhân vật.",
            "• Tab Kỹ năng & Tiềm năng: Nâng cấp kỹ năng chủ động, phân bổ các điểm tiềm năng (Sức mạnh, Khéo léo, Trí tuệ) thu được sau mỗi lần lên cấp để gia tăng sát thương hoặc lượng HP/MP tối đa."
        ],
        "images": ["image_unnamed_64.png", "image_unnamed_65.png"] # Figure 3.19 & 3.20 Chỉ số trang bị & kỹ năng
    },
    {
        "slide_num": 29,
        "title": "Kết quả: Giao diện Chat đa kênh & Tổ đội",
        "subtitle": "Kết quả giao diện cộng đồng",
        "type": "standard",
        "content": [
            "• Chat đa kênh (SignalR): Chat thế giới (World), chat tổ đội (Party), chat riêng tư (Private). Tin nhắn hiển thị lập tức với màu sắc phân biệt kênh rõ ràng.",
            "• Quản lý tổ đội (Party): Người chơi dễ dàng tạo nhóm, gửi lời mời gia nhập, hiển thị thanh máu (HP) của đồng đội thời gian thực trên màn hình để hỗ trợ hồi máu và chiến đấu."
        ],
        "images": ["image_unnamed_67.png", "image_unnamed_69.png"] # Figure 3.22 & 3.24 Chat & Tổ đội
    },
    {
        "slide_num": 30,
        "title": "Kết quả: Giao diện Phó bản Wave & NPC",
        "subtitle": "Kết quả giao diện phó bản",
        "type": "standard",
        "content": [
            "• Giao diện Phó bản Wave: HUD hiển thị số đợt quái vật hiện tại (ví dụ: Wave 3/5), số lượng quái còn lại, và thanh máu của Boss cuối.",
            "• Giao diện tương tác NPC & Nhiệm vụ: Menu động hiển thị danh sách nhiệm vụ từ NPC, cửa hàng bán trang bị, và widget theo dõi mục tiêu nhiệm vụ ở góc màn hình."
        ],
        "images": ["image_unnamed_72.png", "Hình 3.28. Giao diện tương tác NPC nhiệm vụ.png"] # Figure 3.27 & 3.28 Phó bản & NPC
    },
    {
        "slide_num": 31,
        "title": "Kết luận & Định hướng phát triển",
        "subtitle": "Kết luận",
        "type": "standard",
        "content": [
            "• Kết quả đạt được:",
            "  - Hoàn thiện game nhập vai 2D đa người chơi Mutants Arena hoạt động mượt mà, ổn định.",
            "  - Xây dựng thành công cơ chế Gene Ngũ Hành và Hybrid Fusion tạo chiều sâu chiến thuật.",
            "  - Triển khai Dedicated Server Server-Authoritative chống hack/cheat hiệu quả.",
            "• Hạn chế của đề tài: Số lượng bản đồ và quái vật còn ít; giao diện đồ họa 2D pixel art cần được chăm chút bóng bẩy hơn.",
            "• Định hướng tương lai:",
            "  - Phát triển tính năng PvP đấu trường thời gian thực đồng bộ cao.",
            "  - Bổ sung phó bản bang hội, quái vật boss thế giới và tính năng giao dịch vật phẩm giữa người chơi."
        ],
        "images": ["image_unnamed_49.png"] # Figure 3.4 Sơ đồ kiến trúc triển khai
    },
    {
        "slide_num": 32,
        "title": "CẢM ƠN QUÝ THẦY CÔ VÀ CÁC BẠN ĐÃ LẮNG NGHE!",
        "subtitle": "Lời cảm ơn",
        "type": "cover",
        "content": [
            "Sinh viên thực hiện: Trần Văn Thủy",
            "Mã sinh viên: CT060439",
            "Người hướng dẫn khoa học: TS. Nguyễn Đức Hiếu",
            "Khoa Công nghệ thông tin – Học viện Kỹ thuật Mật mã"
        ],
        "images": []
    }
]

# Write Word Document (.docx)
def generate_docx():
    print(f"Creating Word document at {docx_out_path}...")
    doc = docx.Document()
    
    # Title
    title = doc.add_paragraph()
    r = title.add_run("KỊCH BẢN CHI TIẾT VÀ NỘI DUNG SLIDE THUYẾT TRÌNH")
    r.bold = True
    r.font.size = docx.shared.Pt(18)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    subtitle = doc.add_paragraph()
    r = subtitle.add_run("Đề tài: Phát triển trò chơi Mutants Arena với hệ thống tiến hóa Gene bằng Unity")
    r.italic = True
    r.font.size = docx.shared.Pt(12)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph("\n")
    
    for slide in slides_data:
        doc.add_heading(f"Slide {slide['slide_num']}: {slide['title'].replace('\n', ' ')}", level=2)
        
        # Add content as bullet points or text
        for item in slide["content"]:
            p = doc.add_paragraph(style='List Bullet' if item.strip().startswith('•') or item.strip().startswith('-') else 'Normal')
            # Strip bullet marker
            clean_item = item.strip().lstrip('•').lstrip('-').strip()
            p.add_run(clean_item)
            
        # Add images if any
        for img_name in slide["images"]:
            img_path = os.path.join(img_dir, img_name)
            if os.path.exists(img_path):
                try:
                    p = doc.add_paragraph()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p.add_run().add_picture(img_path, width=DocxInches(4.5))
                    p_cap = doc.add_paragraph()
                    p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    r_cap = p_cap.add_run(f"(Hình minh họa nhúng trong slide: {img_name})")
                    r_cap.italic = True
                    r_cap.font.size = docx.shared.Pt(9)
                except Exception as e:
                    print(f"Error adding picture {img_name} to DOCX: {e}")
            else:
                doc.add_paragraph(f"[Cảnh báo: Không tìm thấy ảnh minh họa {img_name}]")
        
        doc.add_paragraph("\n" + "_"*40 + "\n")
        
    doc.save(docx_out_path)
    print("Word document generated successfully!")

# Write clean PowerPoint Document (.pptx) from scratch
def generate_pptx():
    print(f"Creating clean slide deck at {pptx_out_path}...")
    prs = pptx.Presentation()
    
    # Set to standard widescreen (16:9)
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    # Use blank slide layout (usually layout 6 is blank in default template)
    blank_layout = prs.slide_layouts[6]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add slide background color (light clean white-gray)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250)
        
        # 1. Add Title
        title_box = slide.shapes.add_textbox(PptxInches(0.75), PptxInches(0.5), PptxInches(11.833), PptxInches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_data["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(15, 34, 64) # Dark navy blue
        
        # 2. Add Content
        has_imgs = len(slide_data["images"]) > 0
        
        if has_imgs:
            # Split Layout: Text on Left (5.8"), Image on Right (5.8")
            text_width = PptxInches(5.8)
            text_left = PptxInches(0.75)
            
            img_width = PptxInches(5.8)
            img_left = PptxInches(6.8)
        else:
            # Single Column Layout: Full Width (11.833")
            text_width = PptxInches(11.833)
            text_left = PptxInches(0.75)
            
        text_box = slide.shapes.add_textbox(text_left, PptxInches(1.8), text_width, PptxInches(5.0))
        tf_content = text_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = tf_content.margin_top = tf_content.margin_right = tf_content.margin_bottom = 0
        
        for i, item_text in enumerate(slide_data["content"]):
            if i == 0:
                p = tf_content.paragraphs[0]
            else:
                p = tf_content.add_paragraph()
                
            p.text = item_text
            p.font.name = "Arial"
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(50, 50, 50) # Charcoal gray
            p.space_after = Pt(8)
            p.line_spacing = 1.25
            
        # 3. Add images if any
        if has_imgs:
            # If there's 1 image, place it in the center of the right side
            # If there are 2 images, place them vertically stacked or side by side
            if len(slide_data["images"]) == 1:
                img_name = slide_data["images"][0]
                img_path = os.path.join(img_dir, img_name)
                if os.path.exists(img_path):
                    # Scale image to fit inside 5.8" x 5.0" box
                    add_scaled_picture(slide, img_path, img_left, PptxInches(1.8), img_width, PptxInches(5.0))
                else:
                    # Draw a placeholder warning text box
                    warn_box = slide.shapes.add_textbox(img_left, PptxInches(1.8), img_width, PptxInches(5.0))
                    warn_box.text_frame.text = f"[Image {img_name} not found]"
            elif len(slide_data["images"]) >= 2:
                # Two images: stack them vertically or place them side by side
                # Let's stack them vertically with height 2.3" each
                box_height = PptxInches(2.3)
                for j, img_name in enumerate(slide_data["images"][:2]):
                    img_path = os.path.join(img_dir, img_name)
                    top_offset = PptxInches(1.8) + j * PptxInches(2.6)
                    if os.path.exists(img_path):
                        add_scaled_picture(slide, img_path, img_left, top_offset, img_width, box_height)
                    else:
                        warn_box = slide.shapes.add_textbox(img_left, top_offset, img_width, box_height)
                        warn_box.text_frame.text = f"[Image {img_name} not found]"
                        
    prs.save(pptx_out_path)
    alt_pptx_path = pptx_out_path.replace("-clean", "")
    prs.save(alt_pptx_path)
    print(f"PowerPoint presentation generated successfully at {pptx_out_path} and {alt_pptx_path}!")

def generate_md():
    print(f"Creating Markdown slide presentation at {md_out_path}...")
    
    md_content = []
    # Front matter for Markdown Preview Enhanced / Reveal.js presentation mode
    md_content.append("---\npresentation:\n  theme: league\n  width: 1200\n  height: 750\n  margin: 0.1\n---\n\n")
    
    for slide in slides_data:
        title_clean = slide["title"].replace("\n", " ")
        md_content.append(f"# Slide {slide['slide_num']}: {title_clean}\n")
        if slide.get("subtitle"):
            md_content.append(f"## {slide['subtitle']}\n\n")
            
        has_imgs = len(slide["images"]) > 0
        
        # Format bullet points
        bullets = []
        for item in slide["content"]:
            item_clean = item.strip()
            # If it already starts with bullet characters, remove them to normalize
            if item_clean.startswith("•") or item_clean.startswith("-"):
                item_clean = item_clean[1:].strip()
            # Check if there is sub-bullet (e.g. starts with spaces/tab)
            if item.startswith("  -") or item.startswith("  •") or item.startswith("    -"):
                bullets.append(f"  - {item_clean}")
            else:
                bullets.append(f"- {item_clean}")
                
        bullets_str = "\n".join(bullets)
        
        if has_imgs:
            # Multi-column layout using flexbox
            md_content.append("<div style=\"display: flex; gap: 20px; align-items: center;\">\n")
            md_content.append("<div style=\"flex: 1.2;\">\n\n")
            md_content.append(bullets_str)
            md_content.append("\n\n</div>\n")
            md_content.append("<div style=\"flex: 0.8; text-align: center;\">\n\n")
            
            for img_name in slide["images"]:
                # The image folder is c:\Hub\DoAn\extracted_images\
                # Since the md file will be saved in c:\Hub\DoAn\docs\ĐATN-slide.md, 
                # the relative path is ../extracted_images/img_name
                md_content.append(f"![{img_name}](../extracted_images/{img_name})\n\n")
                
            md_content.append("</div>\n</div>\n\n")
        else:
            # Single column
            md_content.append(bullets_str)
            md_content.append("\n\n")
            
        md_content.append("---\n\n")
        
    # Remove the trailing "---"
    if len(md_content) > 0 and md_content[-1] == "---\n\n":
        md_content.pop()
        
    with open(md_out_path, "w", encoding="utf-8") as f:
        f.write("".join(md_content))
    print("Markdown slide presentation generated successfully!")

# Helper function to add scaled picture preserving aspect ratio
def add_scaled_picture(slide, img_path, left_bound, top_bound, max_width, max_height):
    try:
        with PILImage.open(img_path) as img:
            orig_w, orig_h = img.size
            
        aspect_ratio = orig_w / orig_h
        box_ratio = float(max_width) / float(max_height)
        
        if aspect_ratio > box_ratio:
            # Image is wider than target box aspect ratio
            w = max_width
            h = w / aspect_ratio
        else:
            # Image is taller than target box aspect ratio
            h = max_height
            w = h * aspect_ratio
            
        # Center image within the bounding box
        left = left_bound + (max_width - w) / 2
        top = top_bound + (max_height - h) / 2
        
        slide.shapes.add_picture(img_path, left, top, w, h)
        return True
    except Exception as e:
        print(f"Error adding picture {img_path}: {e}")
        return False

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    generate_md()
    print("All file generation completed successfully!")
